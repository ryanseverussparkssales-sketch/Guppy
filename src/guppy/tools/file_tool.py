"""
File reading tool — format-aware text extraction for Guppy.

Supports: PDF, DOCX, XLSX, CSV, PPTX, plain text, code, JSON, YAML, XML, images (EXIF).
All formats are extracted to plain text so the model can read them without
the user manually opening the file.

Token estimation via tiktoken prevents context overflow on large files.
"""
from __future__ import annotations

import csv
import io
import json
import mimetypes
import os
from pathlib import Path
from typing import Any

import chardet
import humanize

# ── Token estimation ──────────────────────────────────────────────────────────

_ENCODER = None


def _count_tokens(text: str) -> int:
    global _ENCODER
    try:
        if _ENCODER is None:
            import tiktoken
            _ENCODER = tiktoken.get_encoding("cl100k_base")
        return len(_ENCODER.encode(text))
    except Exception:
        return len(text) // 4  # rough fallback: ~4 chars per token


def _truncate(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars] + "\n\n[… truncated — file too large for full context]", True


# ── Format readers ────────────────────────────────────────────────────────────

def _read_pdf(path: Path, max_chars: int) -> dict[str, Any]:
    import pdfplumber
    pages_text: list[str] = []
    with pdfplumber.open(path) as pdf:
        total_pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""
            # Also extract tables as plain text
            for table in page.extract_tables():
                rows = ["\t".join(str(c) if c else "" for c in row) for row in table if row]
                if rows:
                    page_text += "\n[Table]\n" + "\n".join(rows)
            if page_text.strip():
                pages_text.append(f"[Page {i + 1}]\n{page_text.strip()}")
            current = "\n\n".join(pages_text)
            if len(current) >= max_chars:
                break
    raw = "\n\n".join(pages_text)
    text, truncated = _truncate(raw, max_chars)
    return {
        "format": "pdf",
        "pages": total_pages,
        "text": text,
        "truncated": truncated,
    }


def _read_docx(path: Path, max_chars: int) -> dict[str, Any]:
    import docx
    doc = docx.Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = para.style.name if para.style else ""
            prefix = "# " if "Heading 1" in style else "## " if "Heading 2" in style else ""
            parts.append(prefix + para.text.strip())
    for table in doc.tables:
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows:
            parts.append("[Table]\n" + "\n".join(rows))
    raw = "\n\n".join(parts)
    text, truncated = _truncate(raw, max_chars)
    return {"format": "docx", "paragraphs": len(doc.paragraphs), "text": text, "truncated": truncated}


def _read_xlsx(path: Path, max_chars: int) -> dict[str, Any]:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append("\t".join("" if c is None else str(c) for c in row))
        if rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    raw = "\n\n".join(parts)
    text, truncated = _truncate(raw, max_chars)
    return {"format": "xlsx", "sheets": len(wb.sheetnames), "text": text, "truncated": truncated}


def _read_csv(path: Path, max_chars: int) -> dict[str, Any]:
    raw_bytes = path.read_bytes()
    enc = chardet.detect(raw_bytes)["encoding"] or "utf-8"
    content = raw_bytes.decode(enc, errors="replace")
    reader = csv.reader(io.StringIO(content))
    rows = ["\t".join(row) for row in reader]
    raw = "\n".join(rows)
    text, truncated = _truncate(raw, max_chars)
    return {"format": "csv", "rows": len(rows), "text": text, "truncated": truncated}


def _read_pptx(path: Path, max_chars: int) -> dict[str, Any]:
    from pptx import Presentation
    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))
    raw = "\n\n".join(parts)
    text, truncated = _truncate(raw, max_chars)
    return {"format": "pptx", "slides": len(prs.slides), "text": text, "truncated": truncated}


def _read_image(path: Path) -> dict[str, Any]:
    from PIL import Image
    from PIL.ExifTags import TAGS
    with Image.open(path) as img:
        info: dict[str, Any] = {
            "format": img.format or path.suffix.lstrip(".").upper(),
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
            "size_px": f"{img.width}×{img.height}",
        }
        # EXIF
        exif_data = getattr(img, "_getexif", lambda: None)()
        if exif_data:
            readable = {}
            for tag_id, val in exif_data.items():
                tag = TAGS.get(tag_id, str(tag_id))
                if isinstance(val, (str, int, float)):
                    readable[tag] = val
            if readable:
                info["exif"] = readable
        lines = [f"{k}: {v}" for k, v in info.items() if k != "exif"]
        if "exif" in info:
            lines.append("EXIF:")
            for k, v in info["exif"].items():
                lines.append(f"  {k}: {v}")
        info["text"] = "\n".join(lines)
        info["truncated"] = False
    return info


def _read_text(path: Path, max_chars: int) -> dict[str, Any]:
    raw_bytes = path.read_bytes()
    enc = chardet.detect(raw_bytes)["encoding"] or "utf-8"
    content = raw_bytes.decode(enc, errors="replace")
    text, truncated = _truncate(content, max_chars)
    return {
        "format": path.suffix.lstrip(".") or "text",
        "encoding": enc,
        "lines": content.count("\n") + 1,
        "text": text,
        "truncated": truncated,
    }


# ── Extension routing ─────────────────────────────────────────────────────────

_BINARY_TEXT_EXTS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv",
    ".pptx", ".ppt",
}
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".heic"}
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".htm", ".css",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env",
    ".xml", ".svg", ".sh", ".bat", ".ps1", ".sql",
    ".c", ".cpp", ".h", ".java", ".go", ".rs", ".rb", ".php",
}


# ── Public API ────────────────────────────────────────────────────────────────

def extract_text(path: str, max_chars: int = 50_000) -> dict[str, Any]:
    """Extract readable text from any supported file format.

    Returns a dict with at minimum: format, text, truncated, token_estimate,
    file_size (human-readable), and format-specific metadata (pages, sheets, etc.).
    """
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(f"File not found: {path}")
    if not p.is_file():
        raise ValueError(f"Path is not a file: {path}")

    stat = p.stat()
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        result = _read_pdf(p, max_chars)
    elif suffix in {".docx", ".doc"}:
        result = _read_docx(p, max_chars)
    elif suffix in {".xlsx", ".xls"}:
        result = _read_xlsx(p, max_chars)
    elif suffix == ".csv":
        result = _read_csv(p, max_chars)
    elif suffix in {".pptx", ".ppt"}:
        result = _read_pptx(p, max_chars)
    elif suffix in _IMAGE_EXTS:
        result = _read_image(p)
    else:
        # Default: try reading as text (covers all code, config, markup, etc.)
        result = _read_text(p, max_chars)

    result["path"] = str(p)
    result["filename"] = p.name
    result["file_size"] = humanize.naturalsize(stat.st_size)
    result["file_size_bytes"] = stat.st_size
    result["token_estimate"] = _count_tokens(result.get("text", ""))
    import arrow
    result["modified"] = arrow.get(stat.st_mtime).humanize()
    return result


def list_directory(path: str, pattern: str = "*") -> dict[str, Any]:
    """List directory contents with human-readable metadata."""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(f"Directory not found: {path}")
    if not p.is_dir():
        raise ValueError(f"Path is not a directory: {path}")

    import arrow
    entries: list[dict[str, Any]] = []
    for child in sorted(p.glob(pattern)):
        stat = child.stat()
        entries.append({
            "name": child.name,
            "type": "dir" if child.is_dir() else "file",
            "size": humanize.naturalsize(stat.st_size) if child.is_file() else None,
            "size_bytes": stat.st_size if child.is_file() else None,
            "modified": arrow.get(stat.st_mtime).humanize(),
            "extension": child.suffix.lower() if child.is_file() else None,
        })
    return {
        "path": str(p),
        "count": len(entries),
        "entries": entries,
    }


def file_info(path: str) -> dict[str, Any]:
    """Return metadata about a file without reading its contents."""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(f"Not found: {path}")
    stat = p.stat()
    import arrow
    return {
        "path": str(p),
        "name": p.name,
        "extension": p.suffix.lower(),
        "type": "dir" if p.is_dir() else "file",
        "size": humanize.naturalsize(stat.st_size),
        "size_bytes": stat.st_size,
        "created": arrow.get(stat.st_ctime).humanize(),
        "modified": arrow.get(stat.st_mtime).humanize(),
        "readable": os.access(p, os.R_OK),
        "writable": os.access(p, os.W_OK),
        "mime_type": mimetypes.guess_type(str(p))[0],
    }
